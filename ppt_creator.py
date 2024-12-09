from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from urllib.request import urlopen
from pptx.oxml import parse_xml

def create_dynamic_ppt(content, heading_font_size, text_font_size, animation, media_url=None):
    prs = Presentation()
    
    animation_effects = {
        "None": None,
        "Fade": "<p:fade />",
        "Slide": "<p:slide />",
        "Zoom": "<p:zoom />"
    }

    for idx, (slide_key, slide_data) in enumerate(content.items()):
        slide_title = slide_data.get('title', '').strip()
        slide_text = slide_data.get('text', '').strip()
        slide_image = slide_data.get('image', 'None').strip()

        if not slide_title or not slide_text:
            continue

        slide_layout = prs.slide_layouts[1]  
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.shapes.placeholders[1]
        
        title.text = slide_title
        content.text = slide_text

        if animation != "None":
            try:
                transition_xml = animation_effects.get(animation, "<p:fade />")
                xml = f'''
                    <mc:AlternateContent xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006">
                        <mc:Choice xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" Requires="p14">
                            <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="slow" p14:dur="3400">
                                {transition_xml}
                            </p:transition>
                        </mc:Choice>
                        <mc:Fallback>
                            <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="slow">
                                <p:fade />
                            </p:transition>
                        </mc:Fallback>
                    </mc:AlternateContent>
                '''
                xmlFragment = parse_xml(xml)
                slide.element.insert(-1, xmlFragment)
            except Exception as e:
                print(f"Error applying transition: {e}")

        if slide_image and slide_image != "None":
            try:
                with urlopen(slide_image) as img_url:
                    img_data = img_url.read()
                image_stream = BytesIO(img_data)
                slide.shapes.add_picture(image_stream, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            except Exception as e:
                print(f"Error adding image: {e}")

        if media_url:
            try:
                add_video_to_slide(media_url, slide)
            except Exception as e:
                print(f"Error adding video: {e}")
    
    return prs

def add_video_to_slide(video_url, slide):
    video_shape = slide.shapes.add_movie(video_url, Inches(0), Inches(0), width=Presentation.slide_width, height=Presentation.slide_height)
    return video_shape
